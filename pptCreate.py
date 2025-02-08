from pptx import Presentation
from llmConnect import generatePPT
from json import loads

def generateSamplePPT():
    # Create a presentation object
    prs = Presentation()    

    slider = prs.slides.add_slide(prs.slide_layouts[0])
    title = slider.shapes.title
    subtitle = slider.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "This is a subtitle"

    # Save the presentation
    prs.save("test.pptx")



    

def generatePPTFromDict(pptDict,file_name="testDataFromLLM.pptx"):
    # Create a presentation object
    prs = Presentation()    

    for slide in pptDict:
         for slide_number, slide_details in slide.items():
            slider = prs.slides.add_slide(prs.slide_layouts[slide_details["slide_layouts"]])

            title = slider.shapes.title
            title.text = slide_details["title"]

            if slide_details["slide_layouts"] == 0:
                subtitle = slider.placeholders[1]
                subtitle.text = slide_details["subtitle"]
            elif slide_details["slide_layouts"] == 1:
                content = slider.placeholders[1]
                content.text = slide_details["subtitle"]

    # Save the presentation
    prs.save(file_name)

if __name__ == "__main__":
    try:
        prompt = "Create a Powerpoint which showcases the VIT University."
        pptDict = loads(generatePPT(prompt))
        print(pptDict)
        print(type(pptDict))
    except Exception as e:
        print("Hit an exception")
        print(e)
    generatePPTFromDict(pptDict)




